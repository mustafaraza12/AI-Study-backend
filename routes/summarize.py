from openai import OpenAI
import os
from dotenv import load_dotenv
import PyPDF2
import docx
from pptx import Presentation
import openpyxl
import re

load_dotenv()

client = OpenAI(
    api_key=os.getenv("GROQ_API_KEY"),
    base_url="https://api.groq.com/openai/v1"
)

# ── Text Extractors ──────────────────────────────────────────────────────────

def extract_from_pdf(path):
    text = ""
    try:
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += (page.extract_text() or "") + "\n"
    except Exception as e:
        print(f"PDF error: {e}")
    return text

def extract_from_word(path):
    text = ""
    try:
        doc = docx.Document(path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Word error: {e}")
    return text

def extract_from_pptx(path):
    text = ""
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            text += f"Slide {i + 1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
            text += "\n"
    except Exception as e:
        print(f"PPT error: {e}")
    return text

def extract_from_excel(path):
    text = ""
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        for sheet in wb.worksheets:
            text += f"Sheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    text += row_text + "\n"
            text += "\n"
    except Exception as e:
        print(f"Excel error: {e}")
    return text

def extract_from_txt(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        print(f"TXT error: {e}")
        return ""

def extract_text(path, filename):
    ext = filename.lower().rsplit(".", 1)[-1]
    if ext == "pdf":
        return extract_from_pdf(path)
    elif ext in ("doc", "docx"):
        return extract_from_word(path)
    elif ext in ("ppt", "pptx"):
        return extract_from_pptx(path)
    elif ext == "xlsx":
        return extract_from_excel(path)
    elif ext == "txt":
        return extract_from_txt(path)
    return ""

# ── Markdown Cleaner ─────────────────────────────────────────────────────────
# Strips all markdown symbols so the frontend receives plain readable text.

def clean_markdown(text: str) -> str:
    # Remove ATX headings (### Title → Title)
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    # Remove bold/italic (**text** / *text* / __text__ / _text_)
    text = re.sub(r"\*{1,3}(.+?)\*{1,3}", r"\1", text)
    text = re.sub(r"_{1,3}(.+?)_{1,3}", r"\1", text)
    # Remove inline code `code`
    text = re.sub(r"`(.+?)`", r"\1", text)
    # Remove horizontal rules
    text = re.sub(r"^[-*_]{3,}\s*$", "", text, flags=re.MULTILINE)
    # Remove blockquotes
    text = re.sub(r"^>\s?", "", text, flags=re.MULTILINE)
    # Collapse 3+ blank lines to 2
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()

# ── Mode Prompts ─────────────────────────────────────────────────────────────

SYSTEM_BASE = (
    "You are a helpful AI study assistant. "
    "CRITICAL FORMATTING RULE: Do NOT use any Markdown syntax in your response. "
    "No hashtags (#), no asterisks (*), no underscores (_), no backticks (`), no hyphens as bullets. "
    "Write in plain text only. Use numbered lists (1. 2. 3.) where lists are needed. "
    "Use blank lines to separate sections. Never hallucinate — if information is missing, say so."
)

MODE_PROMPTS = {
    "summary": {
        "system": SYSTEM_BASE,
        "user": (
            "Read the following document and write a clear, well-structured summary. "
            "Structure your response exactly like this:\n\n"
            "Overview\n"
            "Write 2-3 sentences giving a high-level overview of what the document is about.\n\n"
            "Main Points\n"
            "1. First main point\n"
            "2. Second main point\n"
            "3. (continue as needed)\n\n"
            "Conclusion\n"
            "1-2 sentences on the overall takeaway or conclusion of the document.\n\n"
            "Document:\n{text}"
        ),
    },

    "keypoints": {
        "system": SYSTEM_BASE,
        "user": (
            "Extract the key points from the following document. "
            "Return ONLY a numbered list of the most important ideas, facts, or arguments. "
            "Each point should be one clear sentence. Aim for 5 to 10 points. "
            "Do not add any headers or extra commentary — just the numbered list.\n\n"
            "Document:\n{text}"
        ),
    },

    "flashcards": {
        "system": SYSTEM_BASE,
        "user": (
            "Create 6 to 10 study flashcards from the following document. "
            "Each flashcard must follow this EXACT format with no variation:\n\n"
            "Q: [Question here]\n"
            "A: [Answer here]\n\n"
            "Q: [Question here]\n"
            "A: [Answer here]\n\n"
            "Rules:\n"
            "- Each Q/A pair must be separated by a blank line\n"
            "- Questions should test understanding, not just recall\n"
            "- Answers should be concise (1-2 sentences)\n"
            "- Do not number the cards, do not add headers\n\n"
            "Document:\n{text}"
        ),
    },

    "quiz": {
        "system": SYSTEM_BASE,
        "user": (
            "Create a multiple-choice quiz with 5 questions based on the following document. "
            "Each question must follow this EXACT format with no variation:\n\n"
            "Q1. [Question text]\n"
            "A) [Option]\n"
            "B) [Option]\n"
            "C) [Option]\n"
            "D) [Option]\n"
            "Answer: [Correct letter]\n\n"
            "Q2. [Question text]\n"
            "...\n\n"
            "Rules:\n"
            "- Number questions Q1 through Q5\n"
            "- Always include exactly 4 options A B C D\n"
            "- Always end each question block with 'Answer: X'\n"
            "- Do not add explanations or any other text\n\n"
            "Document:\n{text}"
        ),
    },

    "chat": {
        "system": (
            "You are a helpful AI assistant answering questions about a specific document. "
            "Answer only based on the document content provided. "
            "If the answer is not in the document, say 'That information is not in the document.' "
            "CRITICAL FORMATTING RULE: Do NOT use any Markdown syntax. "
            "No hashtags, asterisks, underscores, or backticks. Plain text only. "
            "Keep answers clear and concise."
        ),
        "user": (
            "Document content:\n{text}\n\n"
            "User question: {question}"
        ),
    },
}

# ── Main Summarizer ──────────────────────────────────────────────────────────

def summarize_text(text: str, mode: str = "summary", question: str = "", max_tokens: int = 1200) -> str:
    """
    Generate AI output for the given mode.

    Modes: summary | keypoints | flashcards | quiz | chat
    """
    mode = mode if mode in MODE_PROMPTS else "summary"
    prompt_config = MODE_PROMPTS[mode]

    # Build the user message
    user_content = prompt_config["user"].format(
        text=text[:10000],   # cap to stay within token limits
        question=question,
    )

    try:
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": prompt_config["system"]},
                {"role": "user",   "content": user_content},
            ],
            max_tokens=max_tokens,
            temperature=0.2,
        )
        raw = response.choices[0].message.content or ""
        # Run the cleaner as a safety net even though we instruct plain text
        return clean_markdown(raw)

    except Exception as e:
        print(f"AI error: {e}")
        return "Error generating output. Please try again."
