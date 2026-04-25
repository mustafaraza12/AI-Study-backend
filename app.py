from flask import Flask, request, jsonify
from flask_cors import CORS
from flask_sqlalchemy import SQLAlchemy
from routes.assignment import solve_assignment
from routes.summarize import extract_text, summarize_text
from routes.quiz import generate_quiz
from routes.code_helper import explain_code
from routes.slide import explain_slide_text
from routes.youtube import fetch_transcript_only, analyze_youtube
from routes.humanize import humanize_text
from routes.math import solve_math, solve_math_image
from routes.essay import write_essay
from routes.flash import generate_flashcards, extract_text_from_file
from flask_jwt_extended import JWTManager
from routes.auth import auth_bp
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from routes.assignmentmaker import write_assignment
from routes.Letter import generate_letter,DOC_LABELS
import traceback
from routes.Iq import generate_questions, evaluate_answers
import logging

import os
import tempfile
from werkzeug.utils import secure_filename
from pptx import Presentation
import PyPDF2
import docx
import openpyxl

app = Flask(__name__)
CORS(app)

# ✅ 1. JWT first
app.config["JWT_SECRET_KEY"] = os.getenv("JWT_SECRET", "aistudysecretkey123")
jwt = JWTManager(app)

# ✅ 2. Register blueprint second
app.register_blueprint(auth_bp, url_prefix="/auth")

# ✅ 3. Limiter third
limiter = Limiter(
    get_remote_address,
    app=app,
    default_limits=["200 per day", "50 per hour"],
    storage_uri="memory://"
)

db_url = os.getenv('DATABASE_URL')

if db_url and db_url.startswith("postgres://"):
    db_url = db_url.replace("postgres://", "postgresql://", 1)

app.config['SQLALCHEMY_DATABASE_URI'] = 'mysql+pymysql://root:Mustafa679/@localhost/studify_db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

# ✅ 4. Apply limits last — after blueprint is registered
limiter.limit("5 per hour")(app.view_functions["auth_routes.register"])
limiter.limit("10 per hour")(app.view_functions["auth_routes.login"])

# Folder for uploaded files
UPLOAD_FOLDER = "./uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

ALLOWED_SLIDE_EXTENSIONS = {"pptx", "ppt", "pdf"}
ALLOWED_DOC_EXTENSIONS   = {"pdf", "doc", "docx", "ppt", "pptx", "xlsx", "txt"}

def allowed_slide_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_SLIDE_EXTENSIONS

def allowed_doc_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_DOC_EXTENSIONS

# ── Text extractors ──────────────────────────────────────────

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    slides_text = []
    for idx, slide in enumerate(prs.slides):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        slides_text.append({"slide_number": idx + 1, "text": text.strip()})
    return slides_text

def extract_text_from_pdf(file_path):
    slides_text = []
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for idx, page in enumerate(reader.pages):
            text = page.extract_text()
            slides_text.append({"slide_number": idx + 1, "text": text.strip() if text else ""})
    return slides_text

def extract_doc_text(file_path, filename):
    """Extract plain text from any supported document type."""
    ext = filename.lower().rsplit(".", 1)[-1]
    text = ""

    if ext == "pdf":
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += (page.extract_text() or "") + "\n"

    elif ext in ("doc", "docx"):
        document = docx.Document(file_path)
        for para in document.paragraphs:
            text += para.text + "\n"

    elif ext in ("ppt", "pptx"):
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text += f"Slide {i + 1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
            text += "\n"

    elif ext == "xlsx":
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet in wb.worksheets:
            text += f"Sheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    text += row_text + "\n"
            text += "\n"

    elif ext == "txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

    return text

# ── Existing Routes ──────────────────────────────────────────

@app.route("/solve-assignment", methods=["POST"])
@limiter.limit("20 per hour")
def solve():
    try:
        data = request.json
        question = data.get("question")
        if not question:
            return jsonify({"error": "No question provided"}), 400
        answer = solve_assignment(question)
        return jsonify({"answer": answer})
    except Exception as e:
        print("SERVER ERROR:", e)
        return jsonify({"error": "Server error"}), 500


# Kept for backward compatibility
@app.route("/summarize-pdf", methods=["POST"])
def summarize_pdf():
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400

        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "No selected file"}), 400

        reader = PyPDF2.PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"

        if not text.strip():
            return jsonify({"error": "PDF is empty"}), 400

        summary = summarize_text(text)
        return jsonify({"summary": summary})

    except Exception as e:
        print("SERVER ERROR (Summarize PDF):", e)
        return jsonify({"error": "Server error"}), 500


@app.route("/summarize-document", methods=["POST"])
@limiter.limit("15 per hour")
def summarize_document():
    try:
        # ── 1. Validate file ────────────────────────────────────────────
        if "file" not in request.files:
            return jsonify({"error": "No file uploaded"}), 400
 
        file = request.files["file"]
        if not file.filename:
            return jsonify({"error": "No file selected"}), 400
 
        if not allowed_doc_file(file.filename):
            return jsonify({"error": "Unsupported file type"}), 400
 
        # ── 2. Read mode + optional chat question from form data ────────
        mode     = request.form.get("mode", "summary").strip().lower()
        question = request.form.get("question", "").strip()   # only used in chat mode
        context  = request.form.get("context", "").strip()    # optional cached summary
 
        VALID_MODES = {"summary", "keypoints", "flashcards", "quiz", "chat"}
        if mode not in VALID_MODES:
            mode = "summary"
 
        # Chat mode requires a question
        if mode == "chat" and not question:
            return jsonify({"error": "Please enter a question to ask about your document."}), 400
 
        # ── 3. Save to temp file ────────────────────────────────────────
        ext = file.filename.rsplit(".", 1)[1].lower()
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
            file.save(tmp.name)
            tmp_path = tmp.name
 
        try:
            # ── 4. Extract text ─────────────────────────────────────────
            text = extract_text(tmp_path, file.filename)
 
            if not text.strip():
                return jsonify({"error": "Could not extract text from this document."}), 422
 
            # ── 5. For chat: prepend cached summary as extra context ────
            # This means the user doesn't need to re-upload the whole file
            # for every follow-up question.
            if mode == "chat" and context:
                text = f"[Document Summary for context]\n{context}\n\n[Full Document Text]\n{text}"
 
            # ── 6. Generate output ──────────────────────────────────────
            result = summarize_text(
                text=text,
                mode=mode,
                question=question,
                max_tokens=1400,
            )
 
            return jsonify({"summary": result})
 
        finally:
            # Always clean up the temp file
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
 
    except Exception as e:
        print("SERVER ERROR (summarize_document):", e)
        return jsonify({"error": "Server error. Please try again."}), 500


 
@app.route("/generate-quiz", methods=["POST"])
@limiter.limit("15 per hour")
def generate_quiz_route():
    try:
        data = request.json
 
        text         = data.get("text", "").strip()
        num_question = int(data.get("num_question", 10))
        difficulty   = data.get("difficulty", "medium").strip().lower()
        quiz_type    = data.get("quiz_type", "mcq").strip().lower()
 
        # ── Validate inputs ──────────────────────────────────────────────────
        if not text:
            return jsonify({"error": "No topic provided"}), 400
 
        if num_question < 1 or num_question > 50:
            return jsonify({"error": "num_question must be between 1 and 50"}), 400
 
        if difficulty not in ("easy", "medium", "hard"):
            difficulty = "medium"
 
        if quiz_type not in ("mcq", "truefalse", "mixed"):
            quiz_type = "mcq"
 
        # ── Generate ─────────────────────────────────────────────────────────
        quiz = generate_quiz(
            text=text,
            num_question=num_question,
            difficulty=difficulty,
            quiz_type=quiz_type,
        )
 
        if not quiz:
            return jsonify({"error": "Could not generate quiz. Try a different topic."}), 400
 
        return jsonify({"quiz": quiz})
 
    except Exception as e:
        print("SERVER ERROR (generate_quiz_route):", e)
        return jsonify({"error": "Server error"}), 500


VALID_EXPLAIN_MODES = {"beginner", "detailed", "summary", "linebyline"}
 
SUPPORTED_LANGUAGES = {
    "python", "javascript", "typescript", "java", "c", "c++", "c#",
    "go", "rust", "php", "ruby", "swift", "kotlin", "sql",
    "html/css", "bash", "dart", "r", "matlab",
}
 
 
@app.route("/explain-code", methods=["POST"])
@limiter.limit("20 per hour")
def explain_code_route():
    data = request.json
 
    if not data:
        return jsonify({"error": "No data provided"}), 400
 
    code = data.get("code", "").strip()
    language = data.get("language", None)
    mode = data.get("mode", "detailed").strip().lower()
 
    if not code:
        return jsonify({"error": "No code provided"}), 400
 
    if len(code) > 8000:
        return jsonify({"error": "Code too long. Please limit to 8000 characters."}), 400
 
    # Validate language
    if language and language.lower() not in SUPPORTED_LANGUAGES:
        language = None  # fallback to auto-detect
 
    # Validate mode
    if mode not in VALID_EXPLAIN_MODES:
        mode = "detailed"
 
    result = explain_code(code, language=language, mode=mode)
 
    if result == "Error explaining code":
        return jsonify({"error": "Failed to explain code. Please try again."}), 500
 
    return jsonify({"explanation": result})


@app.route("/explain-slide", methods=["POST"])
@limiter.limit("10 per hour")
def slide_explainer():
    if "file" not in request.files:
        return jsonify({"explanation": "No file uploaded!"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"explanation": "No file selected!"}), 400

    if not allowed_slide_file(file.filename):
        return jsonify({"explanation": "File type not allowed!"}), 400

    filename = secure_filename(file.filename)
    file_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(file_path)

    ext = filename.rsplit(".", 1)[1].lower()
    if ext in ["pptx", "ppt"]:
        slides = extract_text_from_pptx(file_path)
    elif ext == "pdf":
        slides = extract_text_from_pdf(file_path)
    else:
        return jsonify({"explanation": "Unsupported file type"}), 400

    combined_explanation = ""
    for slide in slides:
        if not slide["text"]:
            continue
        explanation = explain_slide_text(slide["text"], language="Slides")
        combined_explanation += f"### Slide {slide['slide_number']}\n{explanation}\n\n"

    return jsonify({"explanation": combined_explanation})


VALID_MODES = {"standard", "fluency", "formal", "simple", "creative"}
 
@app.route("/humanize-text", methods=["POST"])
@limiter.limit("20 per hour")
def humanize():
    data = request.json
 
    if not data:
        return jsonify({"error": "No data provided"}), 400
 
    text = data.get("text", "").strip()
    mode = data.get("mode", "standard").strip().lower()
 
    if not text:
        return jsonify({"error": "No text provided"}), 400
 
    if len(text) > 5000:
        return jsonify({"error": "Text too long. Please limit to 5000 characters."}), 400
 
    if mode not in VALID_MODES:
        mode = "standard"
 
    result = humanize_text(text, mode)
 
    if result == "Error humanizing text":
        return jsonify({"error": "Failed to humanize text. Please try again."}), 500
 
    return jsonify({"humanized": result})
 
@app.route("/youtube-transcript", methods=["POST"])
@limiter.limit("30 per hour")
def youtube_transcript():
    try:
        data = request.json
        url  = data.get("url", "").strip()
        if not url:
            return jsonify({"error": "No URL provided"}), 400

        result = fetch_transcript_only(url)
        if "error" in result:
            return jsonify(result), 400
        return jsonify(result)

    except Exception as e:
        logging.exception("Transcript endpoint error:")
        return jsonify({"error": "Server error"}), 500


# Slow route — LLM analysis, returns in 8-15s
@app.route("/analyze-youtube", methods=["POST"])
@limiter.limit("10 per hour")
def youtube_analyze():
    try:
        data = request.json
        url  = data.get("url", "").strip()
        if not url:
            return jsonify({"error": "No URL provided"}), 400

        result = analyze_youtube(url)
        if "error" in result:
            return jsonify(result), 400
        return jsonify(result)

    except Exception as e:
        logging.exception("Analyze endpoint error:")
        return jsonify({"error": "Server error"}), 500


@app.route("/solve-math", methods=["POST"])
@limiter.limit("20 per hour")
def math_solver():
    """
    Accepts EITHER:
      A) JSON body  { "problem": "2x + 5 = 11" }
      B) multipart form with:
           - optional "image" file  → OCR'd to text
           - optional "problem" text field
         If both are provided, the extracted OCR text is APPENDED to the
         typed problem so the AI sees everything at once.
    """
    try:
        problem_text = ""
        extracted_ocr = ""
 
        content_type = request.content_type or ""
 
        # ── A: plain JSON ────────────────────────────────────────────────
        if "application/json" in content_type:
            data = request.get_json(force=True) or {}
            problem_text = str(data.get("problem", "")).strip()
 
        # ── B: multipart form (image + optional text) ────────────────────
        else:
            problem_text = str(request.form.get("problem", "")).strip()
 
            if "image" in request.files:
                img_file = request.files["image"]
                if img_file and img_file.filename:
                    from routes.math import extract_text_from_image
                    img_file.seek(0)
                    extracted_ocr = extract_text_from_image(img_file).strip()
 
        # ── Combine typed text + OCR text ────────────────────────────────
        combined = ""
        if problem_text and extracted_ocr:
            combined = f"{problem_text}\n\nProblems from image:\n{extracted_ocr}"
        elif problem_text:
            combined = problem_text
        elif extracted_ocr:
            combined = extracted_ocr
 
        if not combined:
            return jsonify({
                "error": "Please provide a math problem as text or upload an image."
            }), 400
 
        result = solve_math(combined)
 
        if result["success"]:
            response_body = {"data": result["data"]}
            if extracted_ocr:
                response_body["extracted_text"] = extracted_ocr   # show user what was read
            return jsonify(response_body)
        else:
            return jsonify({"error": result.get("error", "Unknown error")}), 500
 
    except Exception as e:
        print("Math Solver Error:", e)
        return jsonify({"error": "Server error"}), 500
 
 
@app.route("/solve-math-image", methods=["POST"])
@limiter.limit("20 per hour")
def math_solver_image():
    """
    Legacy image-only endpoint — kept for backward compatibility.
    Prefer /solve-math with a multipart form now.
    """
    try:
        if "image" not in request.files:
            return jsonify({"error": "No image provided"}), 400
 
        file = request.files["image"]
        result = solve_math_image(file)
 
        if result["success"]:
            return jsonify({
                "data": result["data"],
                "extracted_text": result.get("extracted_text", ""),
            })
        else:
            return jsonify({"error": result.get("error", "Unknown error")}), 500
 
    except Exception as e:
        print("Math Image Solver Error:", e)
        return jsonify({"error": "Server error"}), 500
    

@app.route("/write-essay", methods=["POST"])
@limiter.limit("10 per hour")
def essay_writer():
    try:
        data       = request.json
        topic      = data.get("topic")
        essay_type = data.get("essayType", "argumentative")
        tone       = data.get("tone", "Academic")
        word_count = data.get("wordCount", 500)
        language   = data.get("language", "English")

        if not topic:
            return jsonify({"error": "No topic provided"}), 400

        essay = write_essay(topic, essay_type, tone, word_count, language)
        return jsonify({"essay": essay})

    except Exception as e:
        print("Essay Writer Error:", e)
        return jsonify({"error": "Server error"}), 500
    
@app.route("/generate-flashcards", methods=["POST"])
@limiter.limit("15 per hour")
def flashcard_generator():
    try:
        data       = request.json
        text       = data.get("text", "")
        card_count = int(data.get("cardCount", 10))
        subject    = data.get("subject", "")

        if not text:
            return jsonify({"error": "No text provided"}), 400

        cards = generate_flashcards(text, card_count, subject)

        if not cards:
            return jsonify({"error": "Could not generate flashcards. Try with more content."}), 400

        return jsonify({"flashcards": cards})

    except Exception as e:
        print("Flashcard Error:", e)
        return jsonify({"error": str(e)}), 500


@app.route("/generate-flashcards-file", methods=["POST"])
def flashcard_generator_file():
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file       = request.files["file"]
        card_count = int(request.form.get("cardCount", 10))
        subject    = request.form.get("subject", "")

        text = extract_text_from_file(file, file.filename)

        if not text.strip():
            return jsonify({"error": "Could not extract text from file."}), 422

        cards = generate_flashcards(text, card_count, subject)

        if not cards:
            return jsonify({"error": "Could not generate flashcards from this file."}), 400

        return jsonify({"flashcards": cards})

    except Exception as e:
        print("Flashcard File Error:", e)
        return jsonify({"error": str(e)}), 500


@app.route("/api/assignment", methods=["POST"])
def generate_assignment():
    """
    POST /api/assignment
    Body (JSON):
    {
        "topic"        : "The Impact of Climate Change on Agriculture",
        "pages"        : 3,
        "subject"      : "Environmental Science",
        "student_name" : "Ali Hassan",
        "roll_no"      : "CS-2024-045",
        "professor"    : "Dr. Ahmed Khan",
        "humanize"     : true
    }
 
    Response (JSON):
    {
        "success" : true,
        "content" : "<full assignment body text>"
    }
    """
    try:
        data = request.get_json()
 
        if not data:
            return jsonify({"success": False, "error": "No JSON body provided."}), 400
 
        # ── Required field ───────────────────────────────────────────────────
        topic = data.get("topic", "").strip()
        if not topic:
            return jsonify({"success": False, "error": "Assignment topic is required."}), 400
 
        # ── Optional fields with sensible defaults ───────────────────────────
        pages        = int(data.get("pages", 3))
        subject      = data.get("subject", topic).strip()
        student_name = data.get("student_name", "Student Name").strip()
        roll_no      = data.get("roll_no", "N/A").strip()
        professor    = data.get("professor", "Professor").strip()
        humanize     = bool(data.get("humanize", True))
 
        # ── Validate pages range ─────────────────────────────────────────────
        if pages < 1 or pages > 10:
            return jsonify({"success": False, "error": "Pages must be between 1 and 10."}), 400
 
        # ── Generate ─────────────────────────────────────────────────────────
        content = write_assignment(
            topic        = topic,
            pages        = pages,
            subject      = subject,
            student_name = student_name,
            roll_no      = roll_no,
            professor    = professor,
            humanize     = humanize,
        )
 
        if content.startswith("Error generating assignment:"):
            return jsonify({"success": False, "error": content}), 500
 
        return jsonify({"success": True, "content": content})
 
    except ValueError as ve:
        return jsonify({"success": False, "error": f"Invalid input: {str(ve)}"}), 400
    except Exception as e:
        print(f"[/api/assignment] Unexpected error: {e}")
        return jsonify({"success": False, "error": "An unexpected error occurred."}), 500
 

 
@app.route("/api/health", methods=["GET"])
def health():
    return jsonify({"status": "ok", "service": "IQ Test API"})
 
 
# ── Generate questions ─────────────────────────────────────────────────────
@app.route("/api/iq/generate", methods=["POST"])
def api_generate():
    """
    POST /api/iq/generate
    Body: { name, age, education, subject?, grade?, count? }
    Returns: { success, questions, time_per_question }
    """
    try:
        body = request.get_json(force=True)
 
        name      = str(body.get("name",      "")).strip()
        age       = int(body.get("age",        0))
        education = str(body.get("education", "")).strip()
        subject   = str(body.get("subject",   "")).strip()
        grade     = str(body.get("grade",     "")).strip()
        count     = int(body.get("count",      15))
 
        # Validation
        if not name:
            return jsonify({"success": False, "error": "Name is required"}), 400
        if age < 5 or age > 100:
            return jsonify({"success": False, "error": "Invalid age (must be 5–100)"}), 400
        if not education:
            return jsonify({"success": False, "error": "Education level is required"}), 400
        if count < 5 or count > 30:
            return jsonify({"success": False, "error": "Count must be between 5 and 30"}), 400
 
        questions, time_per_q = generate_questions(
            name=name, age=age, education=education,
            subject=subject, grade=grade, count=count,
        )
 
        return jsonify({
            "success":           True,
            "questions":         questions,
            "time_per_question": time_per_q,
        })
 
    except json.JSONDecodeError as e:
        return jsonify({"success": False, "error": f"AI returned invalid JSON: {str(e)}"}), 500
    except Exception as e:
        traceback.print_exc()
        return jsonify({"success": False, "error": str(e)}), 500
 
 
# ── Evaluate answers ───────────────────────────────────────────────────────
@app.route("/api/iq/evaluate", methods=["POST"])
def api_evaluate():
    """
    POST /api/iq/evaluate
    Body: { name, age, education, subject?, grade?, questions, answers, timed_out, time_taken }
    Returns: { success, result: { iq_score, percentile, correct, accuracy, time_taken,
                                  breakdown, analysis, strengths, improvements } }
    """
    try:
        body = request.get_json(force=True)
 
        name       = str(body.get("name",      "")).strip()
        age        = int(body.get("age",        0))
        education  = str(body.get("education", "")).strip()
        subject    = str(body.get("subject",   "")).strip()
        grade      = str(body.get("grade",     "")).strip()
        questions  = body.get("questions",  [])
        answers    = body.get("answers",    {})
        timed_out  = bool(body.get("timed_out",  False))
        time_taken = int(body.get("time_taken",  0))
 
        if not questions:
            return jsonify({"success": False, "error": "No questions provided"}), 400
 
        result = evaluate_answers(
            name=name, age=age, education=education,
            subject=subject, grade=grade,
            questions=questions, answers=answers,
            timed_out=timed_out, time_taken=time_taken,
        )
 
        return jsonify({"success": True, "result": result})
 
    except json.JSONDecodeError as e:
        return jsonify({"success": False, "error": f"AI returned invalid JSON: {str(e)}"}), 500
    except Exception as e:
        traceback.print_exc()
        return jsonify({"success": False, "error": str(e)}), 500
    

 
 
@app.route("/api/letter", methods=["POST"])
@limiter.limit("15 per hour")
def api_letter():
    try:
        body = request.get_json(force=True)
 
        doc_type        = str(body.get("doc_type",         "")).strip()
        tone            = str(body.get("tone",        "formal")).strip().lower()
        subject         = str(body.get("subject",          "")).strip()
        sender_name     = str(body.get("sender_name",      "")).strip()
        sender_class    = str(body.get("sender_class",     "")).strip()   # ✅ was missing
        recipient_title = str(body.get("recipient_title",  "")).strip()
        recipient_org   = str(body.get("recipient_org",    "")).strip()
        recipient_addr  = str(body.get("recipient_addr",   "")).strip()
        extra_details   = str(body.get("extra_details",    "")).strip()
        date            = str(body.get("date",             "")).strip()   # ✅ was missing
 
        # ── Validation ─────────────────────────────────────────────────────
        if not doc_type:
            return jsonify({"success": False, "error": "doc_type is required"}), 400
 
        if doc_type not in DOC_LABELS:
            return jsonify({
                "success": False,
                "error": f"Unknown doc_type '{doc_type}'.",
            }), 400
 
        if not subject:
            return jsonify({"success": False, "error": "subject is required"}), 400
 
        if not recipient_title:
            return jsonify({"success": False, "error": "recipient_title is required"}), 400
 
        valid_tones = {"formal", "polite", "humble", "urgent"}
        if tone not in valid_tones:
            tone = "formal"
 
        # ── Generate ───────────────────────────────────────────────────────
        content = generate_letter(
            doc_type=doc_type,
            tone=tone,
            subject=subject,
            sender_name=sender_name,
            sender_class=sender_class,       # ✅ correct param name
            recipient_title=recipient_title,
            recipient_org=recipient_org,
            recipient_addr=recipient_addr,
            extra_details=extra_details,
            date=date,                       # ✅ correct param name
        )
 
        return jsonify({"success": True, "content": content})
 
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"success": False, "error": str(e)}), 500
 

db = SQLAlchemy(app)
    

class Post(db.Model):
    __tablename__ = 'Studify' # Tells Flask to look for the 'posts' table
    id = db.Column(db.Integer, primary_key=True)
    title = db.Column(db.String(255), nullable=False)
    slug = db.Column(db.String(255), unique=True, nullable=False)
    content = db.Column(db.Text, nullable=False)
    image_url = db.Column(db.String(500))
    category = db.Column(db.String(50))

# 3. THE ROUTE (The "Slug" logic you mentioned)
@app.route('/api/news/<string:post_slug>', methods=['GET'])
def get_post(post_slug):
    # This searches the DB for a post where the slug matches the URL
    post = Post.query.filter_by(slug=post_slug).first()
    
    if post:
        return jsonify({
            "title": post.title,
            "content": post.content,
            "image": post.image_url,
            "category": post.category
        }), 200
    else:
        return jsonify({"message": "Post not found"}), 404
 
 

if __name__ == "__main__":
    app.run(debug=True, port=5000)
